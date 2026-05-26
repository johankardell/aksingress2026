#!/usr/bin/env python3
"""Build PRESENTATION.pptx from PRESENTATION.md (v2 — polished)."""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
MD   = (ROOT / "PRESENTATION.md").read_text()
OUT  = ROOT / "PRESENTATION.pptx"

# ---------- Theme ----------
BLUE     = RGBColor(0x00, 0x78, 0xD4)
DEEPBLUE = RGBColor(0x10, 0x3A, 0x6E)
DARK     = RGBColor(0x1F, 0x2A, 0x44)
GRAY     = RGBColor(0x60, 0x6E, 0x85)
LIGHT    = RGBColor(0xF3, 0xF5, 0xF9)
SOFT     = RGBColor(0xE6, 0xEE, 0xF7)
ACCENT   = RGBColor(0xE8, 0x1E, 0x4F)
GREEN    = RGBColor(0x2E, 0xA4, 0x4F)
AMBER    = RGBColor(0xE6, 0x97, 0x07)
PURPLE   = RGBColor(0x6E, 0x40, 0xC9)
CODEBG   = RGBColor(0x1E, 0x1E, 0x2E)
CODEFG   = RGBColor(0xE6, 0xE6, 0xE6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# =============================================================================
# Markdown parsing
# =============================================================================
def split_slides(md):
    out, buf, in_code = [], [], False
    for line in md.splitlines():
        if line.startswith("```"):
            in_code = not in_code
            buf.append(line); continue
        if not in_code and line.strip() == "---":
            chunk = "\n".join(buf).strip()
            if chunk: out.append(chunk)
            buf = []; continue
        buf.append(line)
    if buf:
        chunk = "\n".join(buf).strip()
        if chunk: out.append(chunk)
    return out

def parse_slide(block):
    m_part  = re.match(r'^#\s+Part\s+(\d+)\s*[—-]\s*(.+?)$', block, re.MULTILINE)
    m_slide = re.match(r'^##\s+Slide\s+\d+\s*[—-]\s*(.+?)$', block, re.MULTILINE)
    if m_part:
        title = m_part.group(0).lstrip("# ").strip()
        body_md = block[m_part.end():]
        kind = "divider"
    elif m_slide:
        title = m_slide.group(1).strip()
        body_md = block[m_slide.end():]
        kind = "content"
    else:
        title, body_md, kind = "", block, "content"

    speaker = ""
    parts = re.split(r'\*\*Speaker notes[^*]*\*\*\s*', body_md, maxsplit=1)
    if len(parts) == 2:
        body_md, speaker = parts[0], parts[1]

    scripts = []
    def _grab_script(m):
        scripts.append(m.group(1))
        return ""
    body_md = re.sub(
        r'\*\*➡ Move to speaker notes \(script\):\*\*\s*```bash\s*\n(.*?)\n```',
        _grab_script, body_md, flags=re.DOTALL)

    code_blocks = []
    def _grab_code(m):
        code_blocks.append(m.group(1))
        return f"\x00CODE{len(code_blocks)-1}\x00"
    body_md = re.sub(r'```[a-zA-Z0-9]*\n(.*?)\n```', _grab_code, body_md, flags=re.DOTALL)

    notes_full = speaker.strip()
    for sb in scripts:
        notes_full += "\n\n--- Azure CLI script ---\n" + sb.strip()

    return {
        "title": title.replace("Part 1 — ", "").replace("Part 2 — ", "")
                       .replace("Part 3 — ", "").replace("Part 4 — ", "") if kind=="divider" else title,
        "raw_title": title,
        "kind": kind,
        "body_md": body_md.strip(),
        "code_blocks": code_blocks,
        "scripts": scripts,
        "notes": notes_full.strip(),
    }

raw = split_slides(MD)
slide_blocks = [s for s in raw
                if re.search(r'^##\s+Slide\s+\d+', s, re.MULTILINE)
                or re.search(r'^#\s+Part\s+\d+', s, re.MULTILINE)]
parsed = [parse_slide(s) for s in slide_blocks]

# =============================================================================
# Presentation setup
# =============================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# =============================================================================
# Helpers
# =============================================================================
def _i(v):
    return int(v) if v is not None else v

def _add_shape(slide, kind, x, y, w, h):
    return slide.shapes.add_shape(kind, _i(x), _i(y), _i(w), _i(h))

def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _i(x), _i(y), _i(w), _i(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp

def add_round(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _i(x), _i(y), _i(w), _i(h))
    shp.adjustments[0] = 0.10
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=16, bold=False, italic=False,
             color=DARK, font="Calibri", align=PP_ALIGN.LEFT, anchor="top"):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if anchor == "middle":  tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom": tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_line(slide, x1, y1, x2, y2, color=DARK, weight=1.5, arrow_end=True):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _i(x1), _i(y1), _i(x2), _i(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if arrow_end:
        ln = conn.line._get_or_add_ln()
        from pptx.oxml.ns import qn
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return conn

def add_inline_runs(p, text, size=16, color=DARK, font="Calibri"):
    # Bold via **...**, italic via *...*, inline code via `...`
    # First protect inline code, then split bold, then italic.
    tokens = []
    pos = 0
    pattern = re.compile(r'(\*\*[^*]+\*\*)|(`[^`]+`)|(\*[^*\n]+\*)')
    for m in pattern.finditer(text):
        if m.start() > pos:
            tokens.append(("plain", text[pos:m.start()]))
        if m.group(1):   tokens.append(("bold", m.group(1)[2:-2]))
        elif m.group(2): tokens.append(("code", m.group(2)[1:-1]))
        elif m.group(3): tokens.append(("italic", m.group(3)[1:-1]))
        pos = m.end()
    if pos < len(text): tokens.append(("plain", text[pos:]))

    for kind, s in tokens:
        if not s: continue
        r = p.add_run(); r.text = s
        r.font.size = Pt(size)
        if kind == "code":
            r.font.name = "Consolas"
            r.font.color.rgb = RGBColor(0xB0, 0x10, 0x60)
            r.font.bold = False
        else:
            r.font.name = font
            r.font.color.rgb = color
            r.font.bold = (kind == "bold")
            r.font.italic = (kind == "italic")

# ---------- Chrome ----------
def add_title_bar(slide, title, eyebrow=None):
    add_rect(slide, 0, 0, Inches(0.18), Inches(0.95), ACCENT)
    add_rect(slide, Inches(0.18), 0, SW - Inches(0.18), Inches(0.95), DEEPBLUE)
    add_rect(slide, 0, Inches(0.95), SW, Inches(0.04), BLUE)
    if eyebrow:
        add_text(slide, Inches(0.45), Inches(0.10), SW - Inches(0.9), Inches(0.25),
                 eyebrow.upper(), size=10, bold=True,
                 color=RGBColor(0x9C, 0xC8, 0xFF))
        add_text(slide, Inches(0.45), Inches(0.32), SW - Inches(0.9), Inches(0.58),
                 title, size=24, bold=True, color=WHITE, anchor="top")
    else:
        add_text(slide, Inches(0.45), Inches(0.18), SW - Inches(0.9), Inches(0.7),
                 title, size=26, bold=True, color=WHITE, anchor="middle")

def add_footer(slide, page, total, hands_on=False):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), LIGHT)
    add_text(slide, Inches(0.45), SH - Inches(0.30), Inches(8), Inches(0.28),
             "Major Changes in AKS  •  Level 400  •  2026",
             size=9, color=GRAY, anchor="middle")
    if hands_on:
        badge_w = Inches(1.6)
        add_round(slide, SW - Inches(2.0), SH - Inches(0.28), badge_w, Inches(0.24), ACCENT)
        add_text(slide, SW - Inches(2.0), SH - Inches(0.28), badge_w, Inches(0.24),
                 "🛠  HANDS-ON", size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor="middle")
    add_text(slide, SW - Inches(0.8), SH - Inches(0.30), Inches(0.6), Inches(0.28),
             f"{page} / {total}", size=9, bold=True, color=DEEPBLUE,
             align=PP_ALIGN.RIGHT, anchor="middle")

def set_notes(slide, text):
    if not text: return
    nf = slide.notes_slide.notes_text_frame
    nf.clear()
    for i, ln in enumerate(text.split("\n")):
        p = nf.paragraphs[0] if i == 0 else nf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.size = Pt(11)

# =============================================================================
# Renderers for body content (paragraphs, bullets, tables, code/diagrams)
# =============================================================================
def render_text(slide, x, y, w, text, size=14):
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    first = True
    line_count = 0
    for ln in lines:
        if not ln.strip():
            continue
        bullet = False
        indent = 0
        stripped = ln
        m = re.match(r'^(\s*)([-*])\s+(.*)$', ln)
        if m:
            bullet = True
            indent = 1 if len(m.group(1)) >= 2 else 0
            stripped = m.group(3)
        else:
            m2 = re.match(r'^(\s*)(\d+)\.\s+(.*)$', ln)
            if m2:
                bullet = True
                stripped = m2.group(2) + ". " + m2.group(3)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = indent
        p.space_after = Pt(2)
        if bullet:
            r = p.add_run(); r.text = ("▸ " if indent == 0 else "·  ")
            r.font.name = "Calibri"; r.font.size = Pt(size)
            r.font.color.rgb = BLUE; r.font.bold = True
        add_inline_runs(p, stripped, size=size, color=DARK)
        chars = max(1, len(stripped))
        line_count += max(1, (chars // 95) + 1)
    h = Inches(0.06) + Pt(size * 1.4) * line_count
    tb.height = h
    return h

def render_mono_panel(slide, x, y, w, code):
    """Diagram or code in a dark monospaced panel with subtle border."""
    lines = code.split("\n")
    n = len(lines)
    if n > 24:   font_pt, lh = 8, 9.5
    elif n > 16: font_pt, lh = 9, 10.5
    elif n > 10: font_pt, lh = 10, 12.0
    else:        font_pt, lh = 11, 13.0
    h = Inches(0.20) + Pt(lh) * n
    cap_h = Inches(5.0)
    if h > cap_h: h = cap_h
    add_rect(slide, x, y, w, h, CODEBG, line=DEEPBLUE)
    add_rect(slide, x, y, Inches(0.07), h, ACCENT)
    tb = slide.shapes.add_textbox(_i(x + Inches(0.15)), _i(y + Inches(0.08)),
                                  _i(w - Inches(0.25)), _i(h - Inches(0.16)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = "Consolas"; r.font.size = Pt(font_pt)
        r.font.color.rgb = CODEFG
    return h

def render_table(slide, x, y, w, md):
    rows = [r for r in md.split("\n") if r.strip().startswith("|")]
    cells = []
    for r in rows:
        if re.match(r'^\s*\|?\s*[:\-\| ]+\|?\s*$', r): continue
        parts = [c.strip() for c in r.strip().strip("|").split("|")]
        cells.append(parts)
    if not cells: return Inches(0)
    n_rows = len(cells)
    n_cols = max(len(r) for r in cells)
    h = Inches(0.42) * n_rows + Inches(0.05)
    tbl = slide.shapes.add_table(n_rows, n_cols, _i(x), _i(y), _i(w), _i(h)).table
    for j in range(n_cols):
        tbl.columns[j].width = int(w / n_cols)
    for i, row in enumerate(cells):
        for j in range(n_cols):
            cell = tbl.cell(i, j)
            cell.text = ""
            val = row[j] if j < len(row) else ""
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DEEPBLUE
                r = p.add_run(); r.text = val
                r.font.name = "Calibri"; r.font.size = Pt(12)
                r.font.bold = True; r.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
                add_inline_runs(p, val, size=12, color=DARK)
    return h

def render_body(slide, body_md, code_blocks, area):
    x, y, w, h = area
    cursor_y = y
    lines = body_md.splitlines()
    blocks = []
    buf = []
    def flush():
        if buf:
            blocks.append(("text", "\n".join(buf)))
            buf.clear()
    i = 0
    while i < len(lines):
        ln = lines[i]
        m = re.match(r'\x00CODE(\d+)\x00', ln.strip())
        if m:
            flush()
            blocks.append(("code", code_blocks[int(m.group(1))]))
            i += 1; continue
        if ln.strip().startswith("|") and i + 1 < len(lines) and \
                re.match(r'^\s*\|?\s*[:\-\| ]+\|?\s*$', lines[i+1]):
            flush()
            tbl = [ln]; i += 1
            while i < len(lines) and lines[i].strip().startswith("|"):
                tbl.append(lines[i]); i += 1
            blocks.append(("table", "\n".join(tbl)))
            continue
        if ln.strip() == "":
            flush(); i += 1; continue
        buf.append(ln); i += 1
    flush()

    for kind, content in blocks:
        if cursor_y - y > h - Inches(0.3): break
        if kind == "code":
            ch = render_mono_panel(slide, x, cursor_y, w, content)
            cursor_y += ch + Inches(0.10)
        elif kind == "table":
            th = render_table(slide, x, cursor_y, w, content)
            cursor_y += th + Inches(0.15)
        else:
            th = render_text(slide, x, cursor_y, w, content)
            cursor_y += th + Inches(0.06)

# =============================================================================
# Custom slides
# =============================================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DEEPBLUE)
    # decorative diagonal accent
    tri = _add_shape(s, MSO_SHAPE.RIGHT_TRIANGLE, SW - Inches(4.5),
                              SH - Inches(4.5), Inches(4.5), Inches(4.5))
    tri.fill.solid(); tri.fill.fore_color.rgb = BLUE
    tri.line.fill.background(); tri.rotation = 90
    # circle accent
    cir = _add_shape(s, MSO_SHAPE.OVAL, SW - Inches(2.5),
                              Inches(0.6), Inches(1.8), Inches(1.8))
    cir.fill.solid(); cir.fill.fore_color.rgb = ACCENT
    cir.line.fill.background()
    add_text(s, Inches(0.7), Inches(1.2), Inches(3.2), Inches(0.35),
             "AZURE  •  KUBERNETES  •  2026", size=11, bold=True,
             color=RGBColor(0x9C, 0xC8, 0xFF))
    add_text(s, Inches(0.7), Inches(1.7), SW - Inches(1.4), Inches(1.8),
             "Major Changes\nin AKS", size=58, bold=True, color=WHITE)
    add_rect(s, Inches(0.7), Inches(4.4), Inches(0.7), Inches(0.08), ACCENT)
    add_text(s, Inches(0.7), Inches(4.55), SW - Inches(1.4), Inches(0.7),
             "Service Mesh • Ingress • Gateway API • GitOps",
             size=22, color=RGBColor(0xBE, 0xD8, 0xFC))
    add_text(s, Inches(0.7), Inches(5.4), SW - Inches(1.4), Inches(0.5),
             "Level 400  —  deep dive", size=16, italic=True,
             color=RGBColor(0x9C, 0xC8, 0xFF))
    # bottom strip
    add_rect(s, 0, SH - Inches(0.45), SW, Inches(0.45), DARK)
    add_text(s, Inches(0.7), SH - Inches(0.45), SW - Inches(1.4), Inches(0.45),
             "github.com/johankardell/aksingress2026",
             size=11, color=RGBColor(0xBE, 0xD8, 0xFC), anchor="middle")
    return s

def slide_section_divider(part_num, part_title, slide_titles, idx, total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DEEPBLUE)
    add_rect(s, 0, 0, Inches(0.25), SH, ACCENT)
    # huge part number
    add_text(s, Inches(0.7), Inches(0.8), Inches(3), Inches(0.5),
             "PART", size=18, bold=True,
             color=RGBColor(0x9C, 0xC8, 0xFF))
    add_text(s, Inches(0.7), Inches(1.1), Inches(4), Inches(3.5),
             str(part_num), size=240, bold=True, color=WHITE)
    # title
    add_text(s, Inches(5.5), Inches(1.6), SW - Inches(6.2), Inches(1.5),
             part_title, size=40, bold=True, color=WHITE)
    add_rect(s, Inches(5.5), Inches(3.1), Inches(1.0), Inches(0.08), ACCENT)
    # upcoming slide titles
    add_text(s, Inches(5.5), Inches(3.35), Inches(1.8), Inches(0.35),
             "IN THIS PART", size=10, bold=True,
             color=RGBColor(0x9C, 0xC8, 0xFF))
    y = Inches(3.75)
    for t in slide_titles:
        # bullet dot
        dot = _add_shape(s, MSO_SHAPE.OVAL, Inches(5.5), y + Inches(0.10),
                                  Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, Inches(5.8), y, SW - Inches(6.5), Inches(0.42),
                 t, size=15, color=WHITE, anchor="middle")
        y += Inches(0.42)
    add_text(s, Inches(0.7), SH - Inches(0.5), Inches(4), Inches(0.4),
             f"{idx} / {total}",
             size=11, color=RGBColor(0x9C, 0xC8, 0xFF))
    return s

# ---------- Native diagrams ----------
def diag_sidecar(s, x, y, w, h):
    """Two pods each containing an app + envoy sidecar, with mTLS arrow."""
    add_rect(s, x, y, w, h, SOFT)
    add_rect(s, x, y, Inches(0.08), h, BLUE)
    add_text(s, x + Inches(0.2), y + Inches(0.05), w - Inches(0.4), Inches(0.35),
             "Sidecar mode — every pod gets its own Envoy",
             size=12, bold=True, color=DEEPBLUE)

    pod_w = (w - Inches(1.0)) / 2
    pod_h = h - Inches(1.1)
    pod_y = y + Inches(0.55)

    def draw_pod(px, label):
        # pod outer
        outer = _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px, pod_y,
                                    pod_w, pod_h)
        outer.adjustments[0] = 0.08
        outer.fill.solid(); outer.fill.fore_color.rgb = WHITE
        outer.line.color.rgb = DEEPBLUE; outer.line.width = Pt(1.5)
        outer.shadow.inherit = False
        add_text(s, px + Inches(0.1), pod_y + Inches(0.05), pod_w - Inches(0.2),
                 Inches(0.3), label, size=11, bold=True, color=DEEPBLUE)
        # envoy
        envoy_w = (pod_w - Inches(0.45)) / 2
        envoy = add_round(s, px + Inches(0.15), pod_y + Inches(0.45),
                          envoy_w, pod_h - Inches(0.65), AMBER)
        add_text(s, px + Inches(0.15), pod_y + Inches(0.45), envoy_w,
                 pod_h - Inches(0.65), "envoy\nsidecar",
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor="middle")
        # app
        app = add_round(s, px + Inches(0.30) + envoy_w, pod_y + Inches(0.45),
                        envoy_w, pod_h - Inches(0.65), BLUE)
        add_text(s, px + Inches(0.30) + envoy_w, pod_y + Inches(0.45),
                 envoy_w, pod_h - Inches(0.65), "app\ncontainer",
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor="middle")
        return px + envoy_w/2 + Inches(0.15)  # envoy x-center

    cx1 = draw_pod(x + Inches(0.3), "Pod A")
    cx2 = draw_pod(x + Inches(0.7) + pod_w, "Pod B")
    # mTLS arrow between envoys
    arrow_y = pod_y + pod_h/2
    add_line(s, cx1 + Inches(0.45), arrow_y, cx2 - Inches(0.05), arrow_y,
             color=ACCENT, weight=2.2)
    add_text(s, (cx1 + cx2)/2 - Inches(0.6), arrow_y - Inches(0.32),
             Inches(1.2), Inches(0.25), "mTLS", size=10, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)

def diag_ambient(s, x, y, w, h):
    """Two nodes, each with ztunnel daemonset and pods; HBONE tunnel between."""
    add_rect(s, x, y, w, h, SOFT)
    add_rect(s, x, y, Inches(0.08), h, BLUE)
    add_text(s, x + Inches(0.2), y + Inches(0.05), w - Inches(0.4), Inches(0.35),
             "Ambient mode — shared ztunnel, optional waypoint",
             size=12, bold=True, color=DEEPBLUE)

    node_w = (w - Inches(0.9)) / 2
    node_h = h - Inches(0.7)
    node_y = y + Inches(0.5)

    def draw_node(nx, label, with_waypoint=False):
        outer = _add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, nx, node_y,
                                    node_w, node_h)
        outer.adjustments[0] = 0.05
        outer.fill.solid(); outer.fill.fore_color.rgb = WHITE
        outer.line.color.rgb = DEEPBLUE; outer.line.width = Pt(1.5)
        outer.shadow.inherit = False
        add_text(s, nx + Inches(0.1), node_y + Inches(0.05), node_w - Inches(0.2),
                 Inches(0.3), label, size=11, bold=True, color=DEEPBLUE)
        # pods row
        pod_w = (node_w - Inches(0.6)) / 2
        pod_h = Inches(0.85)
        py = node_y + Inches(0.45)
        for j, lbl in enumerate(["app pod", "app pod"]):
            px = nx + Inches(0.2) + j * (pod_w + Inches(0.2))
            p = add_round(s, px, py, pod_w, pod_h, BLUE)
            add_text(s, px, py, pod_w, pod_h, lbl, size=10, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER, anchor="middle")
        # ztunnel DaemonSet bar
        zy = py + pod_h + Inches(0.25)
        zh = Inches(0.55)
        ztun = add_round(s, nx + Inches(0.2), zy, node_w - Inches(0.4), zh, GREEN)
        add_text(s, nx + Inches(0.2), zy, node_w - Inches(0.4), zh,
                 "ztunnel  (DaemonSet — mTLS, identity, L4)",
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor="middle")
        zcx = nx + node_w/2
        zcy = zy + zh/2
        # optional waypoint
        if with_waypoint:
            wy = zy + zh + Inches(0.15)
            wp = add_round(s, nx + Inches(0.6), wy, node_w - Inches(1.2),
                           Inches(0.4), PURPLE)
            add_text(s, nx + Inches(0.6), wy, node_w - Inches(1.2),
                     Inches(0.4), "Waypoint (Envoy, opt-in L7)",
                     size=9, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor="middle")
        return zcx, zcy

    z1x, z1y = draw_node(x + Inches(0.3), "Node 1")
    z2x, z2y = draw_node(x + Inches(0.6) + node_w, "Node 2",
                          with_waypoint=True)
    # HBONE tunnel
    add_line(s, z1x + Inches(0.8), z1y, z2x - Inches(0.8), z2y,
             color=ACCENT, weight=2.5)
    add_text(s, (z1x + z2x)/2 - Inches(1.0), z1y - Inches(0.40),
             Inches(2.0), Inches(0.3),
             "HBONE  •  mTLS  •  HTTP/2 CONNECT  •  :15008",
             size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

def diag_ingress_vs_gateway(s, x, y, w, h):
    """Two columns comparing Ingress (one box) and Gateway API (3 roles)."""
    col_w = (w - Inches(0.4)) / 2
    # Left column — Ingress
    add_rect(s, x, y, col_w, h, SOFT)
    add_rect(s, x, y, col_w, Inches(0.5), GRAY)
    add_text(s, x, y, col_w, Inches(0.5), "Ingress (legacy)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor="middle")
    add_text(s, x + Inches(0.2), y + Inches(0.65), col_w - Inches(0.4),
             Inches(0.35), "One resource, one persona",
             size=11, italic=True, color=GRAY)
    box = add_round(s, x + Inches(0.4), y + Inches(1.15),
                     col_w - Inches(0.8), Inches(1.4), BLUE)
    add_text(s, x + Inches(0.4), y + Inches(1.15), col_w - Inches(0.8),
             Inches(1.4),
             "Ingress\n(host, path, TLS, +\nvendor annotations)",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor="middle")
    add_line(s, x + col_w/2, y + Inches(2.6),
              x + col_w/2, y + Inches(3.0), color=DEEPBLUE, weight=1.5)
    ctl = add_round(s, x + Inches(0.4), y + Inches(3.05),
                     col_w - Inches(0.8), Inches(0.8), DARK)
    add_text(s, x + Inches(0.4), y + Inches(3.05), col_w - Inches(0.8),
             Inches(0.8), "IngressController\n(NGINX, Traefik, …)",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor="middle")
    add_text(s, x + Inches(0.3), y + h - Inches(1.0), col_w - Inches(0.6),
             Inches(0.9),
             "Cluster-scoped object edited by both platform "
             "and app teams.\nPortability via annotation soup.",
             size=11, color=GRAY, italic=True)

    # Right column — Gateway API
    rx = x + col_w + Inches(0.4)
    add_rect(s, rx, y, col_w, h, RGBColor(0xE6, 0xF4, 0xEC))
    add_rect(s, rx, y, col_w, Inches(0.5), GREEN)
    add_text(s, rx, y, col_w, Inches(0.5),
             "Gateway API (current)", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor="middle")
    add_text(s, rx + Inches(0.2), y + Inches(0.65), col_w - Inches(0.4),
             Inches(0.35), "Three resources, three personas",
             size=11, italic=True, color=GREEN)

    # stack: GatewayClass → Gateway → HTTPRoute
    items = [
        ("GatewayClass",  "Infra provider",   DARK),
        ("Gateway",       "Platform team",    DEEPBLUE),
        ("HTTPRoute",     "App developer",    BLUE),
    ]
    box_h = Inches(0.75)
    gap = Inches(0.30)
    bx = rx + Inches(0.3)
    by = y + Inches(1.15)
    for i, (lbl, persona, col) in enumerate(items):
        yb = by + i * (box_h + gap)
        b = add_round(s, bx, yb, col_w - Inches(0.6), box_h, col)
        add_text(s, bx, yb, col_w - Inches(0.6), box_h, lbl,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor="middle")
        add_text(s, bx + col_w - Inches(2.6), yb, Inches(2.0), box_h,
                 persona, size=11, italic=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor="middle")
        if i < len(items) - 1:
            ax = bx + Inches(0.8)
            add_line(s, ax, yb + box_h, ax, yb + box_h + gap,
                     color=GREEN, weight=1.5)
            add_text(s, ax + Inches(0.15), yb + box_h - Inches(0.02),
                     Inches(3.0), gap,
                     "parentRefs ↑" if i == 1 else "GatewayClassName ↑",
                     size=9, italic=True, color=GREEN, anchor="middle")

# ---------- Hands-on (script) slides ----------
def extract_script_steps(script):
    """Pull numbered bash comments (# 1. ...) as steps; fall back to top-level
    # comments. Returns list of (number, text)."""
    steps = []
    for ln in script.split("\n"):
        m = re.match(r'^\s*#\s*(\d+)\.\s+(.+?)\s*$', ln)
        if m:
            steps.append((m.group(1), m.group(2)))
    if not steps:
        for ln in script.split("\n"):
            m = re.match(r'^\s*#\s+(.+?)\s*$', ln)
            if m and len(m.group(1)) > 4 and not m.group(1).startswith("---"):
                steps.append((str(len(steps)+1), m.group(1)))
    return steps

HANDS_ON_INTROS = {
    "Enabling ambient Istio on AKS":
        "Stand up an AKS cluster with the managed Istio add-on in ambient "
        "mode. Workloads get mTLS the moment a namespace is labeled — no "
        "sidecars to inject, no pod restarts on mesh upgrade.",
    "Deploying AGC with Gateway API":
        "Provision Application Gateway for Containers next to an AKS cluster, "
        "wire it in via a delegated subnet, and let the ALB Controller drive "
        "it from standard Gateway API objects using Workload Identity.",
    "Enabling managed Argo CD":
        "Turn AKS into a GitOps target with the managed Argo CD extension. "
        "SSO via Entra ID, target-cluster access via Workload Identity, "
        "and Argo's Application CRDs untouched.",
}
HANDS_ON_OUTCOMES = {
    "Enabling ambient Istio on AKS": [
        "Pinned AKS 1.34.7 with Cilium dataplane in Sweden Central",
        "Managed Istio add-on running, ambient profile active",
        "demo namespace opted in via label — zero pod restarts",
        "Optional namespace-scoped Waypoint for L7 features",
    ],
    "Deploying AGC with Gateway API": [
        "AKS cluster + VNet with a delegated subnet for AGC",
        "User-assigned Managed Identity + federated credential",
        "ALB Controller installed via Helm with Workload Identity",
        "AGC resource + Association live, Gateway + HTTPRoute serving traffic",
    ],
    "Enabling managed Argo CD": [
        "AKS cluster with OIDC issuer and Workload Identity",
        "Managed Argo CD extension installed (Microsoft.ArgoCD)",
        "Entra ID SSO turned on for the Argo CD UI",
        "Bootstrap Application syncing from Git into the demo namespace",
    ],
}

def render_hands_on(s, title, script):
    intro = HANDS_ON_INTROS.get(title, "")
    outcomes = HANDS_ON_OUTCOMES.get(title, [])
    steps = extract_script_steps(script)

    body_y = Inches(1.2)
    body_h = SH - Inches(1.6)

    if intro:
        add_text(s, Inches(0.5), body_y, SW - Inches(1.0), Inches(0.7),
                 intro, size=15, italic=True, color=DEEPBLUE)
        body_y += Inches(0.75)

    # Two columns: left = steps; right = outcomes + script-in-notes callout
    col_gap = Inches(0.3)
    left_w  = Inches(7.5)
    right_x = Inches(0.5) + left_w + col_gap
    right_w = SW - right_x - Inches(0.5)

    # Steps panel
    add_text(s, Inches(0.5), body_y, left_w, Inches(0.35),
             "STEPS PERFORMED BY THE SCRIPT", size=10, bold=True, color=BLUE)
    sy = body_y + Inches(0.40)
    # numbered chips
    for num, text in steps[:8]:
        chip = _add_shape(s, MSO_SHAPE.OVAL, Inches(0.5), sy,
                                   Inches(0.42), Inches(0.42))
        chip.fill.solid(); chip.fill.fore_color.rgb = DEEPBLUE
        chip.line.fill.background()
        add_text(s, Inches(0.5), sy, Inches(0.42), Inches(0.42),
                 num, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor="middle")
        add_text(s, Inches(1.05), sy + Inches(0.02), left_w - Inches(0.6),
                 Inches(0.45), text, size=13, color=DARK, anchor="middle")
        sy += Inches(0.50)

    # Right column: outcomes card
    add_round(s, right_x, body_y, right_w, Inches(2.6), SOFT)
    add_rect(s, right_x, body_y, right_w, Inches(0.45), GREEN)
    add_text(s, right_x, body_y, right_w, Inches(0.45),
             "  ✅  YOU END UP WITH", size=11, bold=True, color=WHITE,
             anchor="middle")
    oy = body_y + Inches(0.55)
    for o in outcomes:
        dot = _add_shape(s, MSO_SHAPE.OVAL, right_x + Inches(0.2),
                                  oy + Inches(0.12), Inches(0.11),
                                  Inches(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()
        add_text(s, right_x + Inches(0.45), oy, right_w - Inches(0.55),
                 Inches(0.45), o, size=12, color=DARK, anchor="middle")
        oy += Inches(0.45)

    # Script callout
    cy = body_y + Inches(2.85)
    add_round(s, right_x, cy, right_w, Inches(1.8), CODEBG, line=ACCENT)
    add_text(s, right_x + Inches(0.25), cy + Inches(0.2), right_w - Inches(0.5),
             Inches(0.35), "📜  FULL AZURE CLI SCRIPT", size=11, bold=True,
             color=ACCENT)
    add_text(s, right_x + Inches(0.25), cy + Inches(0.55),
             right_w - Inches(0.5), Inches(1.2),
             "The complete, copy-pasteable script lives in the\n"
             "speaker notes of this slide.\n\n"
             "Press  Alt+F5  to open Presenter view\n"
             "and copy from the notes pane.",
             size=12, color=CODEFG)

# =============================================================================
# Build the deck
# =============================================================================
# Pre-compute section dividers' upcoming-slide titles
def upcoming_titles(start_idx):
    titles = []
    for p in parsed_iter[start_idx+1:]:
        if p["kind"] == "divider": break
        titles.append(p["title"])
    return titles

def render_agenda(s, body_md):
    """Two-column numbered chip layout for the agenda."""
    # Extract numbered list items
    items = []
    for ln in body_md.splitlines():
        m = re.match(r'^\s*(\d+)\.\s+(.+?)\s*$', ln)
        if m:
            items.append((m.group(1), m.group(2)))
    if not items:
        return
    # Subtitle
    add_text(s, Inches(0.5), Inches(1.15), SW - Inches(1.0), Inches(0.4),
             "Four shifts changing how we build platforms on AKS in 2026",
             size=15, italic=True, color=DEEPBLUE)

    cols = 2
    rows = (len(items) + cols - 1) // cols
    col_w = (SW - Inches(1.4)) / cols
    row_h = Inches(0.78)
    top = Inches(1.85)
    for i, (num, text) in enumerate(items):
        c = i // rows
        r = i % rows
        x = Inches(0.5) + col_w * c
        y = top + row_h * r
        # number chip
        chip = _add_shape(s, MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
        chip.fill.solid()
        chip.fill.fore_color.rgb = ACCENT if int(num) <= 3 else DEEPBLUE
        chip.line.fill.background()
        add_text(s, x, y, Inches(0.6), Inches(0.6), num, size=22, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor="middle")
        # text card
        card = add_round(s, x + Inches(0.75), y + Inches(0.05),
                         col_w - Inches(1.0), Inches(0.5), LIGHT)
        add_text(s, x + Inches(0.9), y + Inches(0.05),
                 col_w - Inches(1.1), Inches(0.5),
                 text, size=14, color=DARK, anchor="middle")
    # bottom hint
    add_text(s, Inches(0.5), SH - Inches(1.0), SW - Inches(1.0), Inches(0.4),
             "→  Q&A throughout; deeper code samples in the speaker notes.",
             size=12, italic=True, color=GRAY)

slide_title_obj = None
def slide_title_capture():
    global slide_title_obj
    slide_title_obj = slide_title()
slide_title_capture()  # Title slide

# The first parsed entry is the "Slide 1 — Title" block; we just rendered the
# title slide specially, so skip it (but preserve its speaker notes).
if parsed and parsed[0]["title"].lower() == "title":
    set_notes(slide_title_obj, parsed[0]["notes"])
    parsed_iter = parsed[1:]
else:
    parsed_iter = parsed


# Eyebrows per part (computed by tracking parts as we walk)
current_part_label = None
current_part_num = 0

# Pre-pass to know real total
preview_total = 1
for p in parsed_iter:
    preview_total += 1
total_slides = preview_total

idx = 1  # slide 1 is the title we already added
for i, p in enumerate(parsed_iter):
    idx += 1
    if p["kind"] == "divider":
        current_part_num += 1
        # Extract part title after the em-dash
        m = re.match(r'Part\s+\d+\s*[—-]\s*(.+)', p["raw_title"])
        ptitle = m.group(1) if m else p["title"]
        current_part_label = ptitle
        slide_section_divider(current_part_num, ptitle,
                              upcoming_titles(i), idx, total_slides)
        continue

    s = prs.slides.add_slide(BLANK)
    short_eyebrow = {
        "How AGC integrates with Gateway API on AKS": "AGC + Gateway API",
    }.get(current_part_label or "", current_part_label or "")
    add_title_bar(s, p["title"], eyebrow=short_eyebrow)

    title_l = p["title"].lower()
    hands_on = bool(p["scripts"])

    # ---- Slide-specific custom rendering ----
    if title_l == "agenda":
        render_agenda(s, p["body_md"])
    elif title_l.startswith("sidecar mode"):
        # Diagram on left, bullets on right
        diag_sidecar(s, Inches(0.5), Inches(1.2), Inches(7.5), Inches(5.2))
        # text from body_md (the bullet list)
        # Strip the first code block (the ASCII diagram we're replacing)
        body_no_code = re.sub(r'\x00CODE\d+\x00', '', p["body_md"]).strip()
        render_text(s, Inches(8.3), Inches(1.4), Inches(4.6), body_no_code,
                    size=14)
    elif title_l.startswith("ambient mode"):
        diag_ambient(s, Inches(0.5), Inches(1.2), Inches(8.3), Inches(5.2))
        body_no_code = re.sub(r'\x00CODE\d+\x00', '', p["body_md"]).strip()
        render_text(s, Inches(9.0), Inches(1.3), Inches(3.9), body_no_code,
                    size=12)
    elif title_l.startswith("ingress vs gateway"):
        diag_ingress_vs_gateway(s, Inches(0.5), Inches(1.2),
                                  SW - Inches(1.0), Inches(5.2))
    elif hands_on:
        render_hands_on(s, p["title"], p["scripts"][0])
    else:
        render_body(s, p["body_md"], p["code_blocks"],
                    area=(Inches(0.5), Inches(1.2),
                          SW - Inches(1.0), SH - Inches(1.6)))

    add_footer(s, idx, total_slides, hands_on=hands_on)
    set_notes(s, p["notes"])

prs.save(OUT)
print(f"Wrote {OUT} with {len(prs.slides)} slides")
